# ingest/loaders.py
from pathlib import Path
from typing import List
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
)
from pptx import Presentation

def _load_pptx(file_path: str) -> List[Document]:
    prs = Presentation(file_path)
    docs: List[Document] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        content = "\n".join(t.strip() for t in texts if t and t.strip())
        if content:
            docs.append(
                Document(
                    page_content=content,
                    metadata={"source": file_path, "slide": i},
                )
            )
    return docs

def load_documents(file_path: str) -> List[Document]:
    """
    Load a single file into a list of Documents with metadata.
    Supported: .pdf, .docx, .pptx/.ppt, .txt, .md
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        loader = PyPDFLoader(file_path)
        return loader.load()

    if ext == ".docx":
        loader = Docx2txtLoader(file_path)
        return loader.load()

    if ext in (".pptx", ".ppt"):
        return _load_pptx(file_path)

    if ext in (".txt", ".md"):
        # TextLoader preserves file path in metadata
        loader = TextLoader(file_path, encoding="utf-8")
        return loader.load()

    raise ValueError(f"Unsupported file type: {ext}")
